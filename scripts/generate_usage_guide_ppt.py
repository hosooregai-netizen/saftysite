from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt


OUT = Path(".artifacts/saftysite-user-guide-controller-worker-2026-04-20.pptx")
LIVE_IMG = Path(".tmp-ui/generated/service-intro")
WORKER_IMG = Path(".artifacts/client-screens/06-worker-calendar.png")
FONT = "Malgun Gothic"
COLORS = {
    "navy": RGBColor(14, 32, 71),
    "teal": RGBColor(19, 122, 127),
    "sand": RGBColor(242, 237, 229),
    "ink": RGBColor(36, 42, 55),
    "muted": RGBColor(95, 103, 122),
    "white": RGBColor(255, 255, 255),
    "coral": RGBColor(226, 119, 92),
    "line": RGBColor(214, 219, 228),
}


def box(slide, x, y, w, h, fill, line=None, radius=True):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        x,
        y,
        w,
        h,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def add_text(slide, x, y, w, h, text, size, color, bold=False):
    frame = slide.shapes.add_textbox(x, y, w, h).text_frame
    frame.word_wrap = True
    for idx, line in enumerate(text.split("\n")):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(2)
    return frame


def add_bullets(slide, x, y, w, h, items, size=19, color=None):
    frame = slide.shapes.add_textbox(x, y, w, h).text_frame
    frame.word_wrap = True
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color or COLORS["ink"]
        p.space_after = Pt(6)
    return frame


def add_title(slide, kicker, title, accent):
    add_text(slide, Inches(0.75), Inches(0.45), Inches(2.5), Inches(0.35), kicker, 15, accent, True)
    add_text(slide, Inches(0.75), Inches(0.78), Inches(11.2), Inches(0.65), title, 28, COLORS["navy"], True)
    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.75), Inches(1.42), Inches(11.0), Inches(0.03)
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = COLORS["line"]
    slide.shapes[-1].line.color.rgb = COLORS["line"]


def add_card(slide, x, y, w, h, title, items, fill):
    box(slide, x, y, w, h, fill, COLORS["line"])
    add_text(slide, x + Inches(0.18), y + Inches(0.15), w - Inches(0.3), Inches(0.35), title, 18, COLORS["navy"], True)
    add_bullets(slide, x + Inches(0.16), y + Inches(0.52), w - Inches(0.28), h - Inches(0.6), items, 15)


def add_screen_slide(prs, kicker, title, image_path, summary, bullets):
    if not image_path.exists():
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box(slide, 0, 0, prs.slide_width, prs.slide_height, COLORS["white"], radius=False)
    add_title(slide, kicker, title, COLORS["teal"])
    box(slide, Inches(0.75), Inches(1.78), Inches(7.55), Inches(4.95), COLORS["sand"], COLORS["line"])
    slide.shapes.add_picture(str(image_path), Inches(0.95), Inches(1.98), width=Inches(7.15), height=Inches(4.55))
    add_card(slide, Inches(8.58), Inches(1.78), Inches(4.0), Inches(4.95), "화면 포인트", [summary, *bullets], COLORS["white"])


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])
box(slide, 0, 0, prs.slide_width, prs.slide_height, COLORS["sand"], radius=False)
box(slide, Inches(0.7), Inches(0.75), Inches(5.0), Inches(5.9), COLORS["navy"])
box(slide, Inches(6.05), Inches(0.95), Inches(6.45), Inches(1.3), COLORS["white"], COLORS["line"])
box(slide, Inches(6.05), Inches(2.45), Inches(3.1), Inches(1.75), COLORS["teal"])
box(slide, Inches(9.4), Inches(2.45), Inches(3.1), Inches(1.75), COLORS["white"], COLORS["line"])
box(slide, Inches(6.05), Inches(4.4), Inches(6.45), Inches(1.75), COLORS["white"], COLORS["line"])
add_text(slide, Inches(1.0), Inches(1.0), Inches(4.2), Inches(0.5), "SAFETY SITE", 18, COLORS["white"], True)
add_text(slide, Inches(1.0), Inches(1.55), Inches(4.2), Inches(1.6), "한국종합안전\n사용 설명서", 30, COLORS["white"], True)
add_text(slide, Inches(1.0), Inches(3.35), Inches(4.1), Inches(1.2), "관제와 지도요원이 실제 업무에서\n어떤 메뉴를 어떻게 쓰는지 빠르게 안내하는 PPT", 18, COLORS["white"])
add_text(slide, Inches(1.0), Inches(5.7), Inches(3.8), Inches(0.5), "2026.04.20", 16, COLORS["white"], True)
add_text(slide, Inches(6.35), Inches(1.2), Inches(5.7), Inches(0.7), "역할별 핵심 메뉴와 사용 흐름 정리", 22, COLORS["navy"], True)
add_text(slide, Inches(6.35), Inches(2.8), Inches(2.4), Inches(0.35), "관제", 20, COLORS["white"], True)
add_text(slide, Inches(6.35), Inches(3.2), Inches(2.35), Inches(0.7), "관리 대시보드\n사업장/현장\n일정·메일·보고서", 16, COLORS["white"])
add_text(slide, Inches(9.7), Inches(2.8), Inches(2.4), Inches(0.35), "지도요원", 20, COLORS["teal"], True)
add_text(slide, Inches(9.7), Inches(3.2), Inches(2.35), Inches(0.7), "배정 현장 확인\n일정 등록·수정\n보고서·사진·메일 발송", 16, COLORS["ink"])
add_text(slide, Inches(6.35), Inches(4.75), Inches(5.7), Inches(0.95), "대상 범위\n관제: 관리 대시보드, 사업장/현장, 일정, 메일함, 실적/매출, 보고서, 사진첩, 콘텐츠\n지도요원: 사업장/현장, 내 일정, 메일함", 16, COLORS["ink"])

slides = [
    ("서비스 개요", "서비스 목적과 역할",
     [("서비스 목적", ["사업장·현장 관리부터 보고서 작성, 발송, 현장 수행까지 하나의 흐름으로 연결합니다.", "관제는 전체 운영을 관리하고, 지도요원은 배정된 현장 업무를 수행합니다."]),
      ("권한 기준", ["관제는 전체 현장과 전체 일정, 전체 매출을 다룹니다.", "지도요원은 자신에게 배정된 현장과 일정만 확인·처리할 수 있습니다."])]),
    ("관제", "관리 대시보드와 실적/매출",
     [("관리 대시보드", ["관리 중 현장 현황을 한눈에 확인합니다.", "교육/계측 자료 충족 상태와 미비 현장을 바로 확인합니다.", "발송 관리 대상, 우선 관리 분기 현황, 종료 임박 현장을 우선순위로 살핍니다."]),
      ("실적 / 매출", ["지도요원 전체 실적과 매출을 요약 지표와 표로 확인합니다.", "기간 기준으로 성과를 비교하고 운영 판단에 활용합니다."])]),
    ("관제", "사업장/현장과 일정",
     [("사업장 / 현장", ["전체 사업장 정보와 현장 기본정보를 확인하고 수정합니다.", "필요 시 지도요원을 현장에 배정하거나 배정을 변경합니다.", "현장 상세 진입 후 보고서, 사진, 이력 화면으로 이동할 수 있습니다."]),
      ("일정 / 캘린더", ["지도요원 전체 일정을 월 단위로 관리합니다.", "현장, 담당자, 상태 기준으로 필터링해 일정 충돌이나 누락을 점검합니다.", "필요한 경우 일정 상태를 수정하고 자료를 내려받을 수 있습니다."])]),
    ("관제", "메일함, 보고서, 사진첩, 콘텐츠",
     [("메일함", ["구글·네이버 메일 계정을 연동해 받은편지함/보낸편지함을 확인합니다.", "보고서를 선택해 메일 작성과 발송을 이어서 처리합니다."]),
      ("보고서", ["전체 보고서를 검색·검토·발송 관리합니다.", "기술지도, 분기, 불량 사업장 관련 문서를 열고 상태를 확인합니다."]),
      ("사진첩 / 콘텐츠", ["사진첩에서는 현장 회차별 사진을 업로드·조회·정리합니다.", "콘텐츠에서는 보고서 작성에 쓰는 공통 자산을 관리합니다."])]),
    ("지도요원", "사업장/현장 사용 방법",
     [("접근 범위", ["지도요원은 배정된 현장 목록만 확인할 수 있습니다.", "배정되지 않은 현장 목록은 확인할 수 없습니다."]),
      ("현장 상세", ["사업장과 현장의 기본 정보를 확인합니다.", "현장별로 기술지도 보고서, 분기 보고서, 불량 사업장, 사진첩 메뉴를 사용합니다.", "업무 중 필요한 문서를 현장 단위로 작성하거나 이어서 수정합니다."])]),
    ("지도요원", "내 일정과 메일함",
     [("내 일정", ["캘린더에서 배정된 현장 일정을 추가합니다.", "선택한 일정의 방문일과 요일, 선택 사유를 수정할 수 있습니다.", "월 단위 달력 보기와 목록 보기를 함께 사용할 수 있습니다."]),
      ("메일함", ["메일 계정을 연동한 뒤 기술지도 보고서를 메일로 발송합니다.", "수신자 주소를 확인하고 발송 이력과 스레드를 함께 확인합니다."])]),
    ("권한 비교", "관제와 지도요원의 차이",
     [("관제 가능", ["전체 현장 조회와 수정", "지도요원 배정 관리", "전체 일정 관리", "전체 매출 관리", "전체 보고서/사진/콘텐츠 관리"]),
      ("지도요원 가능", ["배정된 현장만 조회", "현장 단위 보고서 작성", "내 일정 등록·수정", "현장 사진 관리", "기술지도 보고서 메일 발송"])]),
    ("추천 흐름", "실무 사용 순서",
     [("관제 흐름", ["1. 관리 대시보드로 우선 대응 현장을 확인합니다.", "2. 사업장/현장에서 정보 수정과 배정을 처리합니다.", "3. 일정과 메일함, 보고서에서 후속 조치를 마무리합니다."]),
      ("지도요원 흐름", ["1. 배정된 현장과 내 일정을 확인합니다.", "2. 현장 상세에서 보고서·사진 업무를 수행합니다.", "3. 완료 후 메일함에서 기술지도 보고서를 발송합니다."])]),
]

for kicker, title, groups in slides:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box(slide, 0, 0, prs.slide_width, prs.slide_height, COLORS["white"], radius=False)
    add_title(slide, kicker, title, COLORS["coral"] if kicker in {"관제", "권한 비교"} else COLORS["teal"])
    if len(groups) == 2:
        add_card(slide, Inches(0.75), Inches(1.75), Inches(5.85), Inches(4.9), groups[0][0], groups[0][1], COLORS["sand"])
        add_card(slide, Inches(6.72), Inches(1.75), Inches(5.85), Inches(4.9), groups[1][0], groups[1][1], COLORS["white"])
    elif len(groups) == 3:
        add_card(slide, Inches(0.75), Inches(1.75), Inches(3.8), Inches(4.95), groups[0][0], groups[0][1], COLORS["sand"])
        add_card(slide, Inches(4.77), Inches(1.75), Inches(3.8), Inches(4.95), groups[1][0], groups[1][1], COLORS["white"])
        add_card(slide, Inches(8.79), Inches(1.75), Inches(3.8), Inches(4.95), groups[2][0], groups[2][1], COLORS["sand"])

add_screen_slide(
    prs, "실화면", "관제 관리 대시보드 예시", LIVE_IMG / "overview.png",
    "실데이터 기준으로 KPI 카드와 발송 관리 대상을 바로 확인할 수 있습니다.",
    ["운영 개요와 상태 지표 확인", "지연 보고서 우선 점검", "엑셀 내보내기와 새로고침 지원"],
)
add_screen_slide(
    prs, "실화면", "관제 사업장/현장 관리 예시", LIVE_IMG / "headquarters.png",
    "사업장과 현장을 한 구조로 묶어 검색하고 진입하는 화면입니다.",
    ["사업장 목록 탐색", "현장별 진입 동선 제공", "기본정보 수정과 연결 관리"],
)
add_screen_slide(
    prs, "실화면", "관제 메일함 예시", LIVE_IMG / "mailbox.png",
    "수신 메일, 보낸 메일, 메일 발송 동선을 한 화면에서 관리합니다.",
    ["받은편지함/보낸편지함 구분", "메일 작성과 로그아웃 버튼 제공", "보고서 발송 흐름과 연결"],
)
add_screen_slide(
    prs, "화면 예시", "지도요원 내 일정 예시", WORKER_IMG,
    "지도요원은 월간 일정과 미선택 회차를 같은 흐름에서 확인합니다.",
    ["달력 보기와 목록 보기 사용", "배정 일정만 노출", "방문일과 선택 사유 수정 가능"],
)

slide = prs.slides.add_slide(prs.slide_layouts[6])
box(slide, 0, 0, prs.slide_width, prs.slide_height, COLORS["navy"], radius=False)
add_text(slide, Inches(0.9), Inches(0.95), Inches(5.5), Inches(0.45), "마무리", 18, COLORS["white"], True)
add_text(slide, Inches(0.9), Inches(1.4), Inches(7.0), Inches(1.3), "사용 설명서 핵심 정리", 30, COLORS["white"], True)
add_bullets(slide, Inches(0.95), Inches(2.4), Inches(6.7), Inches(3.2), [
    "관제는 전체 운영 현황과 전체 자원을 관리하는 역할입니다.",
    "지도요원은 배정된 현장 중심으로 일정, 보고서, 사진, 메일 발송을 수행합니다.",
    "두 역할은 같은 데이터를 쓰지만 화면 권한과 접근 범위가 다릅니다.",
], 19, COLORS["white"])
box(slide, Inches(8.0), Inches(1.2), Inches(4.4), Inches(4.9), COLORS["white"])
add_text(slide, Inches(8.3), Inches(1.65), Inches(3.6), Inches(0.55), "권장 안내 문구", 20, COLORS["teal"], True)
add_text(slide, Inches(8.3), Inches(2.3), Inches(3.55), Inches(2.9), "관제는 전체 메뉴를 기준으로 설명하고,\n지도요원은 배정 현장 중심으로\n설명하면 이해가 빠릅니다.\n\n초기 교육 시에는\n관리 대시보드 → 사업장/현장 → 일정 →\n보고서/메일함 순서로 시연하는 것을 권장합니다.", 18, COLORS["ink"])

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(OUT)
